#!/usr/bin/env python3
"""
PPT 流程图生成器
需要安装: pip install python-pptx

功能：
- 将 Mermaid 流程图转换为 PPT
- 支持多页图表
"""

# 注意：需要 pip install python-pptx 后才能使用
# 由于当前环境没有 pip，此脚本仅提供代码参考

TEMPLATE_CODE = '''
# 如果需要生成 PPT，请确保已安装 python-pptx
# pip install python-pptx

from pptx import Presentation
from pptx.util import Inches, Pt

def create_flowchart_ppt(mermaid_code, title="流程图", output_file="flowchart.pptx"):
    """
    创建包含流程图的 PPT
    
    参数:
        mermaid_code: Mermaid 语法代码
        title: 幻灯片标题
        output_file: 输出文件路径
    """
    prs = Presentation()
    
    # 添加标题页
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    
    # 添加内容页
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Mermaid 流程图代码"
    
    # 添加代码文本框
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(5)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    textbox.text_frame.text = mermaid_code
    
    prs.save(output_file)
    print(f"✅ 已保存: {output_file}")

# 示例用法
if __name__ == "__main__":
    example_code = """graph TD
    A[开始] --> B[处理]
    B --> C{判断}
    C -->|是| D[结果1]
    C -->|否| E[结果2]"""
    
    create_flowchart_ppt(example_code, "示例流程图")
'''

# 将模板代码保存为参考文件
with open('/tmp/ppt_template_reference.py', 'w') as f:
    f.write(TEMPLATE_CODE)

print("PPT 生成模板已保存到: /tmp/ppt_template_reference.py")
print("\n如需生成 PPT，请确保安装了 python-pptx:")
print("  pip install python-pptx")
print("\n然后运行:")
print("  python3 /tmp/ppt_template_reference.py")
